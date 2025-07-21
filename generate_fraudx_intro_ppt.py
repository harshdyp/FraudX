from pptx import Presentation
from pptx.util import Pt

slides = [
    {
        "title": "FraudX: Advanced Fraud Detection System",
        "content": [
            "End-to-End Machine Learning Solution",
            "Presented by: [Your Name]"
        ]
    },
    {
        "title": "Project Overview",
        "content": [
            "FraudX is a real-time fraud detection system using machine learning.",
            "Generates and analyzes synthetic transaction data.",
            "Designed for practical deployment and continuous improvement."
        ]
    },
    {
        "title": "System Architecture",
        "content": [
            "Data Generation → Preprocessing → Model Training → Prediction → Results",
            "Modular design: Data, Model, Web App, Batch Processing",
            "Separation of concerns for scalability and maintainability."
        ]
    },
    {
        "title": "Key Features",
        "content": [
            "Realistic synthetic data generation",
            "Explainable AI with feature importance",
            "Interactive Streamlit web app",
            "Automatic model retraining"
        ]
    },
    {
        "title": "What Makes FraudX Unique?",
        "content": [
            "End-to-end, production-ready pipeline",
            "Customizable, realistic data generation",
            "Explainability and transparency",
            "Robust to library upgrades"
        ]
    }
]

def create_presentation(slides, filename="FraudX_Intro_Presentation.pptx"):
    prs = Presentation()
    for slide in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        s = prs.slides.add_slide(slide_layout)
        s.shapes.title.text = slide["title"]
        content = s.placeholders[1]
        content.text = ""
        for bullet in slide["content"]:
            p = content.text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)
        # Remove the first empty paragraph
        if content.text_frame.paragraphs[0].text == "":
            content.text_frame._element.remove(content.text_frame.paragraphs[0]._p)
    prs.save(filename)
    print(f"Presentation saved as {filename}")

if __name__ == "__main__":
    create_presentation(slides) 